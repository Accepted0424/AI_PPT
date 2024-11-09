from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from getPicture import *


# 通过大模型生成的文本制作PPT
def generate(ppt_content):
    ppt = Presentation()
    # PPT首页
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])  # title&subtitle layout
    slide.placeholders[0].text = ppt_content['title']
    slide.placeholders[1].text = "Powered By tongyi.aliyun"

    # 设置首页标题样式
    title_shape = slide.placeholders[0]
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.name = '华文中宋'
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.name = '华文中宋'
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # 内容页
    print('总共%d页...' % len(ppt_content['pages']))
    for i, page in enumerate(ppt_content['pages']):
        print('生成第%d页:%s' % (i + 1, page['title']))
        slide = ppt.slides.add_slide(ppt.slide_layouts[3])  # title&content layout
        slide.placeholders[0].text = page['title']

        # 设置页面标题样式
        slide.placeholders[0].text_frame.paragraphs[0].font.size = Pt(28)
        slide.placeholders[0].text_frame.paragraphs[0].font.bold = True
        slide.placeholders[0].text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 255)

        # 正文
        for sub_content in page['content']:
            if sub_content['title'].strip() and sub_content['description'].strip():
                sub_title = slide.placeholders[1].text_frame.add_paragraph()
                sub_title.text = sub_content['title']
                sub_title.level = 0
                sub_title.font.size = Pt(20)
                sub_title.font.color.rgb = RGBColor(0, 0, 0)

                sub_description = slide.placeholders[1].text_frame.add_paragraph()
                sub_description.text = sub_content['description']
                sub_description.level = 1
                sub_description.font.size = Pt(16)
                sub_description.font.color.rgb = RGBColor(0, 0, 0)

        img_path = getPic("network", "default")
        with Image.open(img_path, "r") as img:
           img_width, img_height = img.size

        img_placeholder = slide.placeholders[2]
        left = img_placeholder.left
        top = img_placeholder.top + Inches(0.5)
        width = Inches(4)
        height = width * (img_height / img_width)

        slide.shapes.add_picture(img_path, left, top, width, height)

        placeholder_to_delete = slide.placeholders[2]
        slide.shapes._spTree.remove(placeholder_to_delete._element)

        # 设置正文对齐方式
        for para in slide.placeholders[1].text_frame.paragraphs:
            para.alignment = PP_ALIGN.LEFT
            para.space_after = Pt(10)

    return ppt
