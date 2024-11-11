import pptx
from pptx.dml.color import RGBColor
from pptx.util import Pt


def generatePPT_with_sample():
    sample_path = "ppt_sample_buaa.pptx"
    ppt = pptx.Presentation(sample_path)

    slide_cover = ppt.slides[0]

    for shape in slide_cover.shapes:
        # 确保该形状包含文本框并且是我们要修改的形状
        if shape.has_text_frame and shape.shape_id == 9:
            title = shape.text_frame.paragraphs[0]
            title.text = "复杂网络的原理"
            title.font.name = "微软雅黑"
            title.font.bold = True
            title.font.size = Pt(66)
            title.font.color.rgb = RGBColor(0,91,171)
            print("已成功修改标题")

    ppt.save('modified_ppt_sample_buaa.pptx')

if __name__ == '__main__':
    generatePPT_with_sample()
