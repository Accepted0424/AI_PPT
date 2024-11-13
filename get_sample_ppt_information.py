from pptx import Presentation

ppt_path = 'ppt_sample_buaa.pptx'
prs = Presentation(ppt_path)

# 遍历每张幻灯片
for slide_num, slide in enumerate(prs.slides, start=1):
    print(f"Slide {slide_num}:")

    # 遍历幻灯片中的所有形状
    for shape_num, shape in enumerate(slide.shapes, start=1):
        if shape.has_text_frame:
            placeholder_text = shape.text if shape.has_text_frame else "<No text>"
            shape_id = shape.shape_id
            print(f"shapeID {shape_id}: {placeholder_text}")
        if shape.shape_type == 13:
            print(f"shape{shape_num} is a image")

    print()  # 每张幻灯片后换行
