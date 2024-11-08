import os
from openai import OpenAI
import json
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import tkinter as tk
from tkinter import messagebox, filedialog
from getPicture import *

pdf_file_path = ''
prompt_file_path = ''
output_file_path = ''


# 用于选择书籍文件的函数
def select_book_file():
    global pdf_file_path
    pdf_file_path = filedialog.askopenfilename(
        title="选择书籍文件",
        filetypes=(("所有文件", "*.*"),)  # 可以根据需要设置文件类型过滤
    )
    if pdf_file_path:
        book_file_label.config(text=f"已选择文件: {pdf_file_path}")


# 用于选择prompt文件的函数
def select_prompt_file():
    global prompt_file_path
    prompt_file_path = filedialog.askopenfilename(
        title="选择prompt文件",
        filetypes=(("文本文件", "*.txt"),)  # 可以根据需要设置文件类型过滤
    )
    if prompt_file_path:
        prompt_file_label.config(text=f"已选择文件: {prompt_file_path}")


def select_output_file():
    global output_file_path
    output_file_path = filedialog.askdirectory(
        title="选择导出位置",
    )
    if output_file_path:
        output_file_label.config(text=f"已选择导出位置: {output_file_path}")


def generate():
    # 获取用户输入
    topic = theme_entry.get()
    pages = page_count_entry.get()

    # 确保页数是一个有效的数字
    try:
        pages = int(pages)
        if pages <= 0:
            raise ValueError
    except ValueError:
        messagebox.showerror("输入错误", "页数必须是正整数！")
        return

    # 创建 OpenAI 客户端
    client = OpenAI(
        api_key="sk-4c1e01470f1d404abbe4eaf23fb3e4d2",
        base_url="https://dashscope.aliyuncs.com/compatible-mode/v1",
    )

    output_format = json.dumps({
        "title": "example title",
        "pages": [
            {
                "title": "title for page 1",
                "content": [
                    {
                        "title": "title for paragraph 1",
                        "description": "detail for paragraph 1",
                    },
                    {
                        "title": "title for paragraph 2",
                        "description": "detail for paragraph 2",
                    },
                ],
            },
            {
                "title": "title for page 2",
                "content": [
                    {
                        "title": "title for paragraph 1",
                        "description": "detail for paragraph 1",
                    },
                    {
                        "title": "title for paragraph 2",
                        "description": "detail for paragraph 2",
                    },
                    {
                        "title": "title for paragraph 3",
                        "description": "detail for paragraph 3",
                    },
                ],
            },
        ],
    }, ensure_ascii=True)

    # 输入 prompt 内容
    prompt = f'''我要准备1个关于{topic}的PPT，要求一共写{pages}页，请你根据主题生成详细内容，不要省略。
    按这个JSON格式输出{output_format}，只能返回JSON，且JSON不要用```包裹，内容要用中文。'''

    # 读取 prompt.txt 文件内容
    with open(prompt_file_path, 'r', encoding='utf-8') as prompt_file:
        prompt_content = prompt_file.read()

    # 读取 book.pdf 文件内容
    pdf_file_path = 'book.pdf'
    pdf_reader = PdfReader(pdf_file_path)
    pdf_content = ""
    for page in pdf_reader.pages:
        pdf_content += page.extract_text()

    # 创建聊天完成请求
    completion = client.chat.completions.create(
        model="qwen-plus",
        messages=[
            {'role': 'system', 'content': 'You are a helpful assistant.'},
            {'role': 'user', 'content': f"文件内容：{pdf_content}"},
            {'role': 'user', 'content': f"内容要求：{prompt_content}"},
            {'role': 'user', 'content': f"格式要求：{prompt}"},
        ]
    )

    # 将结果转换为 JSON 字符串
    completion_json = completion.model_dump_json()
    content = json.loads(completion_json)
    str_content = content['choices'][0]['message']['content']
    str_content = str_content.replace("\\", "\\\\")
    ppt_content = json.loads(str_content)

    ppt = Presentation()

    # PPT首页
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])  # title&subtitle layout
    slide.placeholders[0].text = ppt_content['title']
    slide.placeholders[1].text = "Powered By tongyi.aliyun"

    # 设置首页标题样式
    title_shape = slide.placeholders[0]
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.name = '华文中宋'
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.name = '华文中宋'
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # 内容页
    print('总共%d页...' % len(ppt_content['pages']))
    for i, page in enumerate(ppt_content['pages']):
        print('生成第%d页:%s' % (i + 1, page['title']))
        slide = ppt.slides.add_slide(ppt.slide_layouts[3])  # title&content layout
        slide.placeholders[0].text = page['title']

        # 设置页面标题样式
        slide.placeholders[0].text_frame.paragraphs[0].font.size = Pt(28)
        slide.placeholders[0].text_frame.paragraphs[0].font.bold = True
        slide.placeholders[0].text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 255)

        # 正文
        for sub_content in page['content']:
            if sub_content['title'].strip() and sub_content['description'].strip():
                sub_title = slide.placeholders[1].text_frame.add_paragraph()
                sub_title.text = sub_content['title']
                sub_title.level = 0
                sub_title.font.size = Pt(20)
                sub_title.font.color.rgb = RGBColor(0, 0, 0)

                sub_description = slide.placeholders[1].text_frame.add_paragraph()
                sub_description.text = sub_content['description']
                sub_description.level = 1
                sub_description.font.size = Pt(16)
                sub_description.font.color.rgb = RGBColor(0, 0, 0)


        img_path = getPic("network","default")
        with Image.open(img_path,"r") as img:
            img_width, img_height = img.size

        img_placeholder = slide.placeholders[2]
        left = img_placeholder.left
        top = img_placeholder.top + Inches(0.5)
        width = Inches(4)
        height = width * (img_height/img_width)

        slide.shapes.add_picture(img_path, left, top, width, height)

        placeholder_to_delete = slide.placeholders[2]
        slide.shapes._spTree.remove(placeholder_to_delete._element)

        # 设置正文对齐方式
        for para in slide.placeholders[1].text_frame.paragraphs:
            para.alignment = PP_ALIGN.LEFT
            para.space_after = Pt(10)

    ppt.save(output_file_path + topic + ".ppt")
    output_file_label.config(text=f"PPT导出到: {output_file_path}")
    messagebox.showinfo("导出成功", f"PPT已保存到: {output_file_path}")
    # 关闭GUI窗口
    root.quit()


# 创建GUI窗口
root = tk.Tk()
root.title("PPT生成器")

# 设置窗口大小
root.geometry("500x500")

# 创建输入框和标签
theme_label = tk.Label(root, text="请输入PPT主题:")
theme_label.pack(pady=10)
theme_entry = tk.Entry(root, width=40)
theme_entry.pack(pady=5)

page_count_label = tk.Label(root, text="请输入页数:")
page_count_label.pack(pady=10)
page_count_entry = tk.Entry(root, width=40)
page_count_entry.pack(pady=5)

# 创建“选择书籍文件”按钮
book_button = tk.Button(root, text="选择书籍文件", command=select_book_file)
book_button.pack(pady=10)

# 显示书籍文件路径的标签
book_file_label = tk.Label(root, text="未选择书籍文件")
book_file_label.pack(pady=5)

# 创建“选择Prompt文件”按钮
prompt_button = tk.Button(root, text="选择Prompt文件", command=select_prompt_file)
prompt_button.pack(pady=10)

# 显示Prompt文件路径的标签
prompt_file_label = tk.Label(root, text="未选择Prompt文件")
prompt_file_label.pack(pady=5)

# 创建“选择Prompt文件”按钮
output_button = tk.Button(root, text="导出位置", command=select_output_file)
output_button.pack(pady=10)

# 显示Prompt文件路径的标签
output_file_label = tk.Label(root, text="未选择导出位置")
output_file_label.pack(pady=5)

# 创建生成PPT的按钮
generate_button = tk.Button(root, text="生成PPT", command=generate)
generate_button.pack(pady=20)

# 运行主循环
root.mainloop()